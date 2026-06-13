"""Binary artifact generation (docx / pptx / pdf / xlsx).

The agent gathers the real figures as markdown-ish text; these builders turn
that text into the requested file format under static/files/ and return an
absolute URL. The evaluator scores artifacts on facts + correct format first,
so the priority is that the real data lands in the right container.
"""

from __future__ import annotations

import os
import re
import uuid
from pathlib import Path

from .logging_utils import get_logger

logger = get_logger("artifacts")

_FILES_DIR = Path(__file__).resolve().parent.parent / "static" / "files"

# Explicit format keywords. "slides"/"presentation"/"deck" alone stay inline HTML
# (per AGENTS.md); only an explicit binary format triggers a file.
_FORMAT_PATTERNS: list[tuple[str, str]] = [
    ("pptx", r"\b(pptx|powerpoint|\.ppt)\b"),
    ("docx", r"\b(docx|word document|word doc|\.doc)\b"),
    ("xlsx", r"\b(xlsx|excel|spreadsheet|\.xls|workbook)\b"),
    ("pdf", r"\b(pdf|\.pdf)\b"),
]


def detect_binary_format(question: str) -> str | None:
    q = (question or "").lower()
    for fmt, pattern in _FORMAT_PATTERNS:
        if re.search(pattern, q):
            return fmt
    return None


def public_url(filename: str) -> str:
    base = os.environ.get("PUBLIC_BASE_URL", "http://localhost:8000").rstrip("/")
    return f"{base}/files/{filename}"


def _slug(title: str) -> str:
    slug = re.sub(r"[^a-z0-9]+", "-", (title or "report").lower()).strip("-")
    return slug[:40] or "report"


def _filename(title: str, ext: str) -> str:
    return f"{_slug(title)}-{uuid.uuid4().hex[:8]}.{ext}"


# --- Lightweight markdown parsing --------------------------------------------


def _parse_table(lines: list[str], start: int) -> tuple[list[list[str]], int]:
    """Parse a contiguous markdown table starting at ``start``. Returns
    (rows, next_index). Separator rows (---) are skipped."""
    rows: list[list[str]] = []
    i = start
    while i < len(lines) and "|" in lines[i]:
        cells = [c.strip() for c in lines[i].strip().strip("|").split("|")]
        if not re.fullmatch(r"[\s:|-]+", lines[i].strip()):  # skip --- separators
            rows.append(cells)
        i += 1
    return rows, i


def _blocks(content: str) -> list[dict]:
    """Split content into ordered blocks: heading, table, bullet, paragraph."""
    lines = content.splitlines()
    blocks: list[dict] = []
    i = 0
    while i < len(lines):
        line = lines[i].rstrip()
        if not line.strip():
            i += 1
            continue
        if "|" in line and line.strip().startswith("|"):
            rows, i = _parse_table(lines, i)
            if rows:
                blocks.append({"type": "table", "rows": rows})
            continue
        heading = re.match(r"^(#{1,6})\s+(.*)$", line)
        if heading:
            blocks.append({"type": "heading", "level": len(heading.group(1)), "text": heading.group(2)})
            i += 1
            continue
        bullet = re.match(r"^[-*]\s+(.*)$", line.strip())
        if bullet:
            blocks.append({"type": "bullet", "text": bullet.group(1)})
            i += 1
            continue
        blocks.append({"type": "paragraph", "text": line.strip()})
        i += 1
    return blocks


# --- Builders -----------------------------------------------------------------


def _build_docx(title: str, content: str, path: Path) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading(title, level=0)
    for b in _blocks(content):
        if b["type"] == "heading":
            doc.add_heading(b["text"], level=min(b["level"], 4))
        elif b["type"] == "bullet":
            doc.add_paragraph(b["text"], style="List Bullet")
        elif b["type"] == "table":
            rows = b["rows"]
            cols = max(len(r) for r in rows)
            table = doc.add_table(rows=0, cols=cols)
            table.style = "Light Grid Accent 1"
            for r in rows:
                cells = table.add_row().cells
                for j, val in enumerate(r):
                    cells[j].text = val
        else:
            doc.add_paragraph(b["text"])
    doc.save(str(path))


def _build_pdf(title: str, content: str, path: Path) -> None:
    from fpdf import FPDF
    from fpdf.enums import XPos, YPos

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    def text(s: str) -> str:  # fpdf core fonts are latin-1 only
        return s.encode("latin-1", "replace").decode("latin-1")

    def line(s: str, height: float = 6.0) -> None:
        # new_x=LMARGIN/new_y=NEXT resets the cursor to the left margin so the
        # next multi_cell has the full page width available.
        pdf.multi_cell(0, height, text(s), new_x=XPos.LMARGIN, new_y=YPos.NEXT)

    pdf.set_font("Helvetica", "B", 18)
    line(title, 10)
    pdf.ln(2)
    for b in _blocks(content):
        if b["type"] == "heading":
            pdf.set_font("Helvetica", "B", max(16 - b["level"], 11))
            line(b["text"], 8)
        elif b["type"] == "bullet":
            pdf.set_font("Helvetica", "", 11)
            line(f"- {b['text']}")
        elif b["type"] == "table":
            pdf.set_font("Helvetica", "", 10)
            for r in b["rows"]:
                line(" | ".join(r))
        else:
            pdf.set_font("Helvetica", "", 11)
            line(b["text"])
        pdf.ln(1)
    pdf.output(str(path))


def _build_pptx(title: str, content: str, path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title

    # Group blocks into slides by heading; bullets/paragraphs become body lines.
    sections: list[tuple[str, list[str]]] = []
    current: tuple[str, list[str]] | None = None
    for b in _blocks(content):
        if b["type"] == "heading":
            if current:
                sections.append(current)
            current = (b["text"], [])
        else:
            if current is None:
                current = (title, [])
            if b["type"] == "table":
                for r in b["rows"]:
                    current[1].append(" | ".join(r))
            else:
                current[1].append(b["text"])
    if current:
        sections.append(current)
    if not sections:
        sections = [(title, [content.strip() or "(no content)"])]

    for heading, lines in sections:
        # Chunk long sections across multiple slides.
        for chunk_start in range(0, max(len(lines), 1), 8):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = heading
            body = slide.placeholders[1].text_frame
            body.clear()
            chunk = lines[chunk_start : chunk_start + 8] or [""]
            for k, line in enumerate(chunk):
                para = body.paragraphs[0] if k == 0 else body.add_paragraph()
                para.text = line
                para.font.size = Pt(16)
    prs.save(str(path))


def _build_xlsx(title: str, content: str, path: Path) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = (title or "Sheet")[:31]
    row = 1
    wrote_table = False
    for b in _blocks(content):
        if b["type"] == "table":
            wrote_table = True
            for r in b["rows"]:
                for col, val in enumerate(r, start=1):
                    ws.cell(row=row, column=col, value=val)
                row += 1
            row += 1
        elif b["type"] == "heading":
            ws.cell(row=row, column=1, value=b["text"])
            row += 1
        elif not wrote_table:
            ws.cell(row=row, column=1, value=b["text"])
            row += 1
    if row == 1:
        ws.cell(row=1, column=1, value=content.strip() or "(no content)")
    wb.save(str(path))


_BUILDERS = {
    "docx": _build_docx,
    "pdf": _build_pdf,
    "pptx": _build_pptx,
    "xlsx": _build_xlsx,
}


def build_artifact(fmt: str, title: str, content: str) -> str:
    """Build the requested artifact and return its public URL.

    Raises ValueError for an unknown format.
    """
    if fmt not in _BUILDERS:
        raise ValueError(f"unsupported artifact format: {fmt}")
    _FILES_DIR.mkdir(parents=True, exist_ok=True)
    filename = _filename(title, fmt)
    _BUILDERS[fmt](title, content, _FILES_DIR / filename)
    url = public_url(filename)
    logger.info("Built %s artifact -> %s", fmt, url)
    return url
